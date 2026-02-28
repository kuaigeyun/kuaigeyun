#!/usr/bin/env python3
"""
从 README 生成快格轻制造产品介绍 PPT
包含全部 23 张功能截图，专业配色、版式与装饰
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 路径
SCRIPT_DIR = Path(__file__).resolve().parent
PROJECT_ROOT = SCRIPT_DIR.parent
SCREENSHOTS_DIR = PROJECT_ROOT / "screenshots"
OUTPUT_PATH = PROJECT_ROOT / "快格轻制造-产品介绍.pptx"

# 品牌配色
COLORS = {
    "primary": RGBColor(0, 51, 102),      # 深蓝 #003366
    "secondary": RGBColor(0, 102, 161),    # 中蓝 #0066A1
    "accent": RGBColor(0, 163, 224),       # 亮蓝 #00A3E0
    "light": RGBColor(245, 248, 250),      # 浅灰 #F5F8FA
    "text": RGBColor(51, 51, 51),         # 正文 #333333
    "text_light": RGBColor(102, 102, 102),
    "white": RGBColor(255, 255, 255),
}

# 每页 1 张截图 + 1 个设计亮点
SCREENSHOT_ITEMS = [
    ("1. 登录页面.png", "登录页面", "简洁登录，支持多租户与主题切换"),
    ("2. 工作台.png", "工作台", "数据概览一目了然，关键指标集中展示"),
    ("3. 主题配置.png", "主题配置", "支持亮色/暗色主题，满足不同使用场景"),
    ("4. 插件式应用.png", "插件式应用", "模块按需开启，轻量起步、灵活扩展"),
    ("5. 业务蓝图配置.png", "业务蓝图配置", "可视化流程编排，业务节点可配置"),
    ("6. 自定义字段设置.png", "自定义字段设置", "满足企业个性化字段扩展需求"),
    ("7. 外部数据连接.png", "外部数据连接", "对接外部数据源，实现数据集成"),
    ("8. 外部应用连接.png", "外部应用连接", "与第三方系统集成，打通企业上下游"),
    ("9. 审批流设计.png", "审批流设计", "图形化审批流编排，流程可配置、可追溯"),
    ("10. 定时消息.png", "定时消息", "定时消息推送，重要事项不错过"),
    ("11. 打印模板设计.png", "打印模板设计", "自由设计打印模板，支持报表定制"),
    ("12. BI大屏设计.png", "BI 大屏设计", "数据可视化大屏，经营看板一目了然"),
    ("13. 智能建议.png", "智能建议", "AI 辅助决策，智能建议辅助操作"),
    ("14. 图形化BOM设计.png", "图形化 BOM 设计", "可视化 BOM 编辑，结构清晰易维护"),
    ("15. 工单便捷操作.png", "工单便捷操作", "工单快捷操作，一键下推、批量处理"),
    ("16. 触屏终端界面.png", "触屏终端界面", "工位机大按钮、大字体，扫码报工、数字键盘录入"),
    ("17. 手机端集成.png", "手机端集成", "移动端适配，支持现场与远程协同"),
    ("18. 锁屏界面.png", "锁屏界面", "工位机锁屏保护，防止误操作"),
    ("19. 全流程跟踪.png", "全流程跟踪", "需求→工单→入库全链路追溯，变更影响可分析"),
    ("20. 高级搜索面板.png", "高级搜索面板", "多条件筛选，快速定位目标数据"),
    ("21. 甘特图排产.png", "甘特图排产", "工单级排产可视化，时间轴一目了然"),
    ("22. 嵌入式EXCEL在线导入.png", "嵌入式 Excel 在线导入", "Univer 嵌入式 Excel，在线导入导出降低录入门槛"),
    ("23. 上线检查助手.png", "上线检查助手", "引导式配置检查，确保上线前配置完整"),
]


def set_slide_background(slide, color):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, top=0, height=0.05):
    """顶部装饰条"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(top), Inches(10), Inches(height)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["secondary"]
    bar.line.fill.background()


def add_title_slide(prs, title, subtitle=""):
    """封面页：渐变感背景 + 居中标题"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS["primary"])
    # 顶部装饰条
    add_accent_bar(slide, 0, 0.08)
    # 主标题
    left, top, width, height = Inches(0.5), Inches(2.2), Inches(9), Inches(1.3)
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]
    p.alignment = PP_ALIGN.CENTER
    # 副标题
    if subtitle:
        left, top, width, height = Inches(0.5), Inches(3.6), Inches(9), Inches(1)
        tf2 = slide.shapes.add_textbox(left, top, width, height).text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(22)
        p2.font.color.rgb = RGBColor(180, 210, 235)
        p2.alignment = PP_ALIGN.CENTER
    # 底部装饰
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.8), Inches(10), Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0, 40, 80)
    bar.line.fill.background()
    return slide


def add_content_slide(prs, title, bullets):
    """内容页：浅色背景 + 左侧色条 + 标题 + 要点"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS["light"])
    add_accent_bar(slide, 0, 0.06)
    # 左侧标题色条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(0.08), Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["secondary"]
    bar.line.fill.background()
    # 标题
    left, top, width, height = Inches(0.7), Inches(0.45), Inches(8.5), Inches(0.75)
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    # 内容区（带浅色底）
    content_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.35), Inches(9), Inches(5.5))
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = COLORS["white"]
    content_box.line.color.rgb = RGBColor(220, 230, 240)
    content_box.shadow.inherit = False
    # 内容文字
    left, top, width, height = Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.2)
    tf2 = slide.shapes.add_textbox(left, top, width, height).text_frame
    tf2.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS["text"]
        p.space_after = Pt(10)
    return slide


def add_table_slide(prs, title, rows):
    """表格页：带表头样式的表格"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS["light"])
    add_accent_bar(slide, 0, 0.06)
    # 左侧色条 + 标题
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(0.08), Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["secondary"]
    bar.line.fill.background()
    left, top, width, height = Inches(0.7), Inches(0.45), Inches(8.5), Inches(0.6)
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    # 表格
    num_rows, num_cols = len(rows), len(rows[0])
    left, top = Inches(0.5), Inches(1.2)
    table = slide.shapes.add_table(num_rows, num_cols, left, top, Inches(9), Inches(0.35 * num_rows)).table
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(cell_text)[:60]
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["secondary"]
                para.font.bold = True
                para.font.color.rgb = COLORS["white"]
            else:
                para.font.color.rgb = COLORS["text"]
    return slide


def add_single_feature_slide(prs, title, image_path, design_highlight):
    """功能页：图片带边框 + 设计亮点卡片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS["light"])
    add_accent_bar(slide, 0, 0.05)
    # 标题区（带色条）
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(0.06), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["accent"]
    bar.line.fill.background()
    left, top, width, height = Inches(0.65), Inches(0.35), Inches(8.5), Inches(0.55)
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    # 图片容器（圆角矩形底 + 图片）
    img_left, img_top = Inches(0.6), Inches(1.05)
    img_width, img_height = Inches(8.8), Inches(4.2)
    # 白色底框（模拟图片边框）
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, img_left, img_top, img_width, img_height)
    frame.fill.solid()
    frame.fill.fore_color.rgb = COLORS["white"]
    frame.line.color.rgb = RGBColor(200, 215, 230)
    frame.line.width = Pt(1)
    # 图片（略小于框，留白边）
    pic_left = img_left + Inches(0.05)
    pic_top = img_top + Inches(0.05)
    pic_width = img_width - Inches(0.1)
    pic_height = img_height - Inches(0.1)
    slide.shapes.add_picture(str(image_path), pic_left, pic_top, width=pic_width, height=pic_height)
    # 设计亮点卡片
    card_top = Inches(5.4)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), card_top, Inches(9), Inches(1.5))
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["secondary"]
    card.line.fill.background()
    left, top, width, height = Inches(0.7), Inches(5.55), Inches(8.6), Inches(1.2)
    tf2 = slide.shapes.add_textbox(left, top, width, height).text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"✦ 设计亮点：{design_highlight}"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = COLORS["white"]
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, "快格轻制造", "面向中小制造企业的轻量级 MES · 产品介绍")

    add_content_slide(prs, "产品简介 · 核心定位", [
        "轻量起步：专注生产执行核心环节，小企业可极简运行",
        "灵活扩展：按需开启销售、采购、质量、财务等模块",
        "需求驱动：销售预测与销售订单统一为「需求」，一套流程覆盖多种业务",
        "面向中国：针对一人多岗、数据录入量大、快速上手等痛点优化",
        "",
        "典型用户：10–200 人离散制造企业，多品种小批量、按单或按库生产",
    ])

    add_table_slide(prs, "技术栈", [
        ["层级", "技术"],
        ["前端", "React 18 + TypeScript + Vite、Ant Design Pro"],
        ["后端", "FastAPI、Tortoise ORM + PostgreSQL"],
        ["架构", "多租户 SaaS、插件化应用"],
    ])

    add_table_slide(prs, "已上线模块", [
        ["模块", "功能概览"],
        ["销售管理", "报价单、试样打样、销售预测、销售订单、发货通知"],
        ["计划管理", "需求管理、需求计算、生产计划、排程管理、计算配置"],
        ["采购管理", "采购申请、采购订单、到货通知"],
        ["生产执行", "工单、报工、工位机触屏、返工单、委外、异常管理"],
        ["质量管理", "来料/过程/成品检验、检验计划、质量追溯"],
        ["设备/工装", "设备、模具、工装、故障、保养计划、保养提醒"],
        ["仓储管理", "入库/出库、线边仓、倒冲、库存、盘点、调拨、组装/拆解"],
        ["分析中心", "单据时效、成本核算、全追溯、经营看板"],
        ["财务管理", "应收、应付、发票、收款、付款"],
        ["绩效管理", "节假日、技能、员工配置、计件/工时、KPI、绩效汇总"],
    ])

    add_content_slide(prs, "业务流程 · 统一主线", [
        "需求来源（销售预测或销售订单）→ 需求计算 → 工单/采购 → 生产执行 → 入库/出库",
        "",
        "销售订单 → 下推生成 Demand → 需求计算 → 工单/采购",
        "",
        "计划四模块：需求管理 → 需求计算 → 生产计划 → 排程管理",
        "",
        "典型流程：MTO/MTS 统一主线；线边仓报工触发 BOM 倒冲",
    ])

    add_content_slide(prs, "产品优势", [
        "功能设计：统一需求计算引擎、双路径可配置、线边仓+报工倒冲、全链路文档追溯",
        "",
        "交互设计：列表+Modal+Drawer、工位机触屏大按钮、批量与一键、生命周期可视化",
        "",
        "痛点解决：轻量设计、减少步骤、批量导入、统一主线、流程开关配置",
    ])

    for filename, title, design_highlight in SCREENSHOT_ITEMS:
        path = SCREENSHOTS_DIR / filename
        if path.exists():
            add_single_feature_slide(prs, title, path, design_highlight)
        else:
            print(f"警告：截图不存在 {path}")

    add_content_slide(prs, "快速开始", [
        "环境：Node.js 22、Python 3.12、PostgreSQL 15、Redis 6",
        "",
        "git clone <repository-url>",
        "cd riveredge",
        "./Launch.dev.sh",
    ])

    add_content_slide(prs, "联系方式", [
        "官网：https://kuaigeyun.com",
        "邮箱：ludingjie@live.cn",
        "V: lu_dingjie",
    ])

    prs.save(str(OUTPUT_PATH))
    print(f"✅ PPT 已生成：{OUTPUT_PATH}")


if __name__ == "__main__":
    main()
